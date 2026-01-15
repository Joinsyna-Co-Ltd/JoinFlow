"""
Advanced Export Functionality
=============================

Enterprise-grade export capabilities including Excel, PowerPoint, and enhanced formats.
Multi-format output capabilities for enterprise use.
"""

import io
import json
import logging
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import html
import re

logger = logging.getLogger(__name__)


# ============================================
# Excel Exporter
# ============================================

class ExcelExporter:
    """Excel报告导出器 - 企业级功能"""
    
    @staticmethod
    def is_available() -> bool:
        """检查Excel导出是否可用"""
        try:
            import openpyxl
            return True
        except ImportError:
            return False
    
    @staticmethod
    def export_task_result(
        task_id: str,
        description: str,
        result: str,
        steps: List[Dict],
        metadata: Dict = None
    ) -> bytes:
        """导出任务结果为Excel"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.chart import BarChart, Reference
            
            wb = Workbook()
            
            # ===== 概览页 =====
            ws_overview = wb.active
            ws_overview.title = "任务概览"
            
            # 样式定义
            header_font = Font(bold=True, size=14, color="FFFFFF")
            header_fill = PatternFill(start_color="4361EE", end_color="4361EE", fill_type="solid")
            title_font = Font(bold=True, size=16)
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # 标题
            ws_overview['A1'] = "📋 JoinFlow 任务执行报告"
            ws_overview['A1'].font = Font(bold=True, size=20, color="4361EE")
            ws_overview.merge_cells('A1:D1')
            
            # 基本信息
            info_data = [
                ("任务ID", task_id),
                ("任务描述", description[:100] + "..." if len(description) > 100 else description),
                ("创建时间", metadata.get('created_at', datetime.now().isoformat()) if metadata else datetime.now().isoformat()),
                ("状态", metadata.get('status', 'completed') if metadata else 'completed'),
                ("总步骤数", str(len(steps))),
                ("完成步骤", str(sum(1 for s in steps if s.get('status') == 'completed'))),
            ]
            
            row = 3
            for label, value in info_data:
                ws_overview[f'A{row}'] = label
                ws_overview[f'A{row}'].font = Font(bold=True)
                ws_overview[f'B{row}'] = value
                row += 1
            
            # 调整列宽
            ws_overview.column_dimensions['A'].width = 15
            ws_overview.column_dimensions['B'].width = 60
            
            # ===== 执行步骤页 =====
            ws_steps = wb.create_sheet("执行步骤")
            
            # 表头
            headers = ["序号", "步骤名称", "Agent", "状态", "输出摘要"]
            for col, header in enumerate(headers, 1):
                cell = ws_steps.cell(row=1, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
                cell.border = border
            
            # 数据
            for i, step in enumerate(steps, 1):
                status = step.get('status', 'pending')
                status_text = {
                    'completed': '✅ 完成',
                    'failed': '❌ 失败',
                    'running': '🔄 执行中',
                    'pending': '⏳ 待执行'
                }.get(status, status)
                
                output = step.get('output', '')
                output_summary = output[:100] + "..." if len(output) > 100 else output
                
                row_data = [
                    i,
                    step.get('description', ''),
                    step.get('agent', 'unknown'),
                    status_text,
                    output_summary
                ]
                
                for col, value in enumerate(row_data, 1):
                    cell = ws_steps.cell(row=i+1, column=col, value=value)
                    cell.border = border
                    if col == 4:  # 状态列颜色
                        if 'completed' in status:
                            cell.fill = PatternFill(start_color="D4EDDA", end_color="D4EDDA", fill_type="solid")
                        elif 'failed' in status:
                            cell.fill = PatternFill(start_color="F8D7DA", end_color="F8D7DA", fill_type="solid")
            
            # 调整列宽
            ws_steps.column_dimensions['A'].width = 8
            ws_steps.column_dimensions['B'].width = 30
            ws_steps.column_dimensions['C'].width = 15
            ws_steps.column_dimensions['D'].width = 12
            ws_steps.column_dimensions['E'].width = 50
            
            # ===== 结果页 =====
            ws_result = wb.create_sheet("执行结果")
            ws_result['A1'] = "最终结果"
            ws_result['A1'].font = title_font
            ws_result['A3'] = result
            ws_result['A3'].alignment = Alignment(wrap_text=True)
            ws_result.column_dimensions['A'].width = 100
            
            # ===== 统计页 =====
            ws_stats = wb.create_sheet("统计分析")
            
            # 状态统计
            status_counts = {}
            for step in steps:
                status = step.get('status', 'pending')
                status_counts[status] = status_counts.get(status, 0) + 1
            
            ws_stats['A1'] = "执行状态统计"
            ws_stats['A1'].font = title_font
            
            ws_stats['A3'] = "状态"
            ws_stats['B3'] = "数量"
            ws_stats['A3'].font = header_font
            ws_stats['B3'].font = header_font
            ws_stats['A3'].fill = header_fill
            ws_stats['B3'].fill = header_fill
            
            row = 4
            for status, count in status_counts.items():
                ws_stats[f'A{row}'] = status
                ws_stats[f'B{row}'] = count
                row += 1
            
            # 添加图表
            if len(status_counts) > 0:
                chart = BarChart()
                chart.title = "步骤状态分布"
                chart.type = "col"
                data = Reference(ws_stats, min_col=2, min_row=3, max_row=3+len(status_counts))
                cats = Reference(ws_stats, min_col=1, min_row=4, max_row=3+len(status_counts))
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws_stats.add_chart(chart, "D3")
            
            # 保存到内存
            buffer = io.BytesIO()
            wb.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            logger.warning("openpyxl not installed")
            raise ImportError("需要安装 openpyxl: pip install openpyxl")
        except Exception as e:
            logger.error(f"Excel export failed: {e}")
            raise
    
    @staticmethod
    def export_data_table(
        data: List[Dict],
        title: str = "数据报告",
        sheet_name: str = "数据"
    ) -> bytes:
        """将数据列表导出为Excel表格"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            if not data:
                ws['A1'] = "无数据"
                buffer = io.BytesIO()
                wb.save(buffer)
                return buffer.getvalue()
            
            # 样式
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="4361EE", end_color="4361EE", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # 标题
            ws['A1'] = title
            ws['A1'].font = Font(bold=True, size=16)
            
            # 表头
            headers = list(data[0].keys())
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=3, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            # 数据
            for row_idx, row_data in enumerate(data, 4):
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ''))
                    cell.border = border
            
            # 自动调整列宽
            for col_idx, header in enumerate(headers, 1):
                max_length = len(str(header))
                for row_data in data:
                    value = str(row_data.get(header, ''))
                    if len(value) > max_length:
                        max_length = min(len(value), 50)
                ws.column_dimensions[chr(64 + col_idx)].width = max_length + 2
            
            buffer = io.BytesIO()
            wb.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            raise ImportError("需要安装 openpyxl: pip install openpyxl")


# ============================================
# PowerPoint Exporter
# ============================================

class PowerPointExporter:
    """PowerPoint演示文稿导出器 - 企业级功能"""
    
    @staticmethod
    def is_available() -> bool:
        """检查PPT导出是否可用"""
        try:
            from pptx import Presentation
            return True
        except ImportError:
            return False
    
    @staticmethod
    def export_task_result(
        task_id: str,
        description: str,
        result: str,
        steps: List[Dict],
        metadata: Dict = None
    ) -> bytes:
        """导出任务结果为PowerPoint"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor as RgbColor  # 兼容新版本
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # 颜色定义
            primary_color = RgbColor(67, 97, 238)  # #4361EE
            success_color = RgbColor(46, 204, 113)  # #2ECC71
            danger_color = RgbColor(231, 76, 60)   # #E74C3C
            dark_color = RgbColor(26, 26, 46)      # #1A1A2E
            
            # ===== 标题页 =====
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 背景
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = dark_color
            background.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "📋 任务执行报告"
            title_para.font.size = Pt(48)
            title_para.font.bold = True
            title_para.font.color.rgb = RgbColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
            
            # 副标题
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = description[:80] + "..." if len(description) > 80 else description
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RgbColor(148, 163, 184)
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            # 日期
            date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.5))
            date_frame = date_box.text_frame
            date_para = date_frame.paragraphs[0]
            date_para.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            date_para.font.size = Pt(14)
            date_para.font.color.rgb = RgbColor(100, 116, 139)
            date_para.alignment = PP_ALIGN.CENTER
            
            # ===== 任务概览页 =====
            slide = prs.slides.add_slide(slide_layout)
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "任务概览"
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = primary_color
            
            # 信息卡片
            info_items = [
                ("🆔 任务ID", task_id),
                ("📝 描述", description[:60] + "..." if len(description) > 60 else description),
                ("📊 状态", metadata.get('status', 'completed') if metadata else 'completed'),
                ("📈 完成率", f"{sum(1 for s in steps if s.get('status') == 'completed')}/{len(steps)} 步骤"),
            ]
            
            y_pos = Inches(1.3)
            for label, value in info_items:
                # 标签
                label_box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(3), Inches(0.5))
                label_frame = label_box.text_frame
                label_para = label_frame.paragraphs[0]
                label_para.text = label
                label_para.font.size = Pt(18)
                label_para.font.bold = True
                
                # 值
                value_box = slide.shapes.add_textbox(Inches(4), y_pos, Inches(8), Inches(0.5))
                value_frame = value_box.text_frame
                value_para = value_frame.paragraphs[0]
                value_para.text = str(value)
                value_para.font.size = Pt(18)
                
                y_pos += Inches(0.8)
            
            # ===== 执行步骤页 =====
            # 每页显示3个步骤
            steps_per_page = 3
            for page_idx in range(0, len(steps), steps_per_page):
                page_steps = steps[page_idx:page_idx + steps_per_page]
                
                slide = prs.slides.add_slide(slide_layout)
                
                # 标题
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"执行步骤 ({page_idx + 1}-{min(page_idx + steps_per_page, len(steps))}/{len(steps)})"
                title_para.font.size = Pt(36)
                title_para.font.bold = True
                title_para.font.color.rgb = primary_color
                
                y_pos = Inches(1.3)
                for i, step in enumerate(page_steps):
                    step_idx = page_idx + i + 1
                    status = step.get('status', 'pending')
                    
                    # 状态图标和颜色
                    if status == 'completed':
                        icon = "✅"
                        status_color = success_color
                    elif status == 'failed':
                        icon = "❌"
                        status_color = danger_color
                    else:
                        icon = "⏳"
                        status_color = RgbColor(100, 116, 139)
                    
                    # 步骤卡片背景
                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), y_pos, Inches(12), Inches(1.8)
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = RgbColor(248, 250, 252)
                    card.line.color.rgb = status_color
                    
                    # 步骤标题
                    step_title = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.2), Inches(11), Inches(0.5))
                    step_frame = step_title.text_frame
                    step_para = step_frame.paragraphs[0]
                    step_para.text = f"{icon} 步骤 {step_idx}: {step.get('description', '')[:50]}"
                    step_para.font.size = Pt(20)
                    step_para.font.bold = True
                    
                    # Agent信息
                    agent_box = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.7), Inches(11), Inches(0.4))
                    agent_frame = agent_box.text_frame
                    agent_para = agent_frame.paragraphs[0]
                    agent_para.text = f"Agent: {step.get('agent', 'unknown')} | 状态: {status}"
                    agent_para.font.size = Pt(14)
                    agent_para.font.color.rgb = RgbColor(100, 116, 139)
                    
                    # 输出摘要
                    output = step.get('output', '')
                    if output:
                        output_box = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(1.1), Inches(11), Inches(0.5))
                        output_frame = output_box.text_frame
                        output_para = output_frame.paragraphs[0]
                        output_para.text = output[:100] + "..." if len(output) > 100 else output
                        output_para.font.size = Pt(12)
                        output_para.font.color.rgb = RgbColor(71, 85, 105)
                    
                    y_pos += Inches(2)
            
            # ===== 结果页 =====
            slide = prs.slides.add_slide(slide_layout)
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "🎯 执行结果"
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = primary_color
            
            # 结果内容
            result_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
            result_frame = result_box.text_frame
            result_frame.word_wrap = True
            result_para = result_frame.paragraphs[0]
            # 截取结果前1500字符
            result_text = result[:1500] + "..." if len(result) > 1500 else result
            result_para.text = result_text
            result_para.font.size = Pt(14)
            
            # ===== 结束页 =====
            slide = prs.slides.add_slide(slide_layout)
            
            # 背景
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = dark_color
            background.line.fill.background()
            
            # 感谢文字
            thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
            thanks_frame = thanks_box.text_frame
            thanks_para = thanks_frame.paragraphs[0]
            thanks_para.text = "✨ 报告完成"
            thanks_para.font.size = Pt(48)
            thanks_para.font.bold = True
            thanks_para.font.color.rgb = RgbColor(255, 255, 255)
            thanks_para.alignment = PP_ALIGN.CENTER
            
            # 品牌
            brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.5))
            brand_frame = brand_box.text_frame
            brand_para = brand_frame.paragraphs[0]
            brand_para.text = "Powered by JoinFlow"
            brand_para.font.size = Pt(18)
            brand_para.font.color.rgb = RgbColor(100, 116, 139)
            brand_para.alignment = PP_ALIGN.CENTER
            
            # 保存到内存
            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            logger.warning("python-pptx not installed")
            raise ImportError("需要安装 python-pptx: pip install python-pptx")
        except Exception as e:
            logger.error(f"PPT export failed: {e}")
            raise


# ============================================
# Enhanced Export Manager
# ============================================

class AdvancedExportManager:
    """高级导出管理器 - 企业级功能"""
    
    def __init__(self, output_dir: str = "./exports"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def export_task(
        self,
        task_id: str,
        description: str,
        result: str,
        steps: List[Dict],
        format: str = "markdown",
        metadata: Dict = None,
        save_to_file: bool = True
    ) -> Tuple[Any, str]:
        """
        导出任务结果到多种格式
        
        Args:
            format: markdown, html, json, pdf, excel, pptx
            save_to_file: 是否保存到文件
            
        Returns:
            (content, file_path) 元组
        """
        from .exporter import MarkdownExporter, HTMLExporter, JSONExporter, PDFExporter
        
        format_lower = format.lower()
        content = None
        ext = format_lower
        is_binary = False
        
        if format_lower in ('markdown', 'md'):
            content = MarkdownExporter.export_task_result(task_id, description, result, steps, metadata)
            ext = 'md'
            
        elif format_lower == 'html':
            content = HTMLExporter.export_task_result(task_id, description, result, steps, metadata)
            
        elif format_lower == 'json':
            content = JSONExporter.export_task_result(task_id, description, result, steps, metadata)
            
        elif format_lower == 'pdf':
            if not PDFExporter.is_available():
                raise ImportError("PDF导出需要安装 reportlab: pip install reportlab")
            content = PDFExporter.export_task_result(task_id, description, result, steps, metadata)
            is_binary = True
            
        elif format_lower in ('excel', 'xlsx'):
            if not ExcelExporter.is_available():
                raise ImportError("Excel导出需要安装 openpyxl: pip install openpyxl")
            content = ExcelExporter.export_task_result(task_id, description, result, steps, metadata)
            ext = 'xlsx'
            is_binary = True
            
        elif format_lower in ('pptx', 'ppt', 'powerpoint'):
            if not PowerPointExporter.is_available():
                raise ImportError("PPT导出需要安装 python-pptx: pip install python-pptx")
            content = PowerPointExporter.export_task_result(task_id, description, result, steps, metadata)
            ext = 'pptx'
            is_binary = True
            
        else:
            raise ValueError(f"不支持的导出格式: {format}")
        
        file_path = ""
        if save_to_file:
            # 生成安全的文件名
            safe_desc = re.sub(r'[\\/*?:"<>|]', '_', description[:30])
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{safe_desc}_{timestamp}.{ext}"
            file_path = str(self.output_dir / filename)
            
            if is_binary:
                with open(file_path, 'wb') as f:
                    f.write(content)
            else:
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(content)
        
        return content, file_path
    
    def get_available_formats(self) -> List[Dict]:
        """获取所有可用的导出格式"""
        from .exporter import PDFExporter
        
        formats = [
            {
                "id": "markdown",
                "name": "Markdown",
                "extension": ".md",
                "icon": "📝",
                "available": True,
                "description": "轻量级标记语言，适合文档和笔记"
            },
            {
                "id": "html",
                "name": "HTML",
                "extension": ".html",
                "icon": "🌐",
                "available": True,
                "description": "网页格式，可直接在浏览器查看"
            },
            {
                "id": "json",
                "name": "JSON",
                "extension": ".json",
                "icon": "📊",
                "available": True,
                "description": "结构化数据格式，便于程序处理"
            },
            {
                "id": "pdf",
                "name": "PDF",
                "extension": ".pdf",
                "icon": "📄",
                "available": PDFExporter.is_available(),
                "description": "便携文档格式，适合打印和分享"
            },
            {
                "id": "excel",
                "name": "Excel",
                "extension": ".xlsx",
                "icon": "📈",
                "available": ExcelExporter.is_available(),
                "description": "电子表格格式，支持数据分析和图表"
            },
            {
                "id": "pptx",
                "name": "PowerPoint",
                "extension": ".pptx",
                "icon": "📽️",
                "available": PowerPointExporter.is_available(),
                "description": "演示文稿格式，适合汇报展示"
            },
        ]
        return formats
    
    def export_batch(
        self,
        tasks: List[Dict],
        formats: List[str] = None
    ) -> Dict[str, List[str]]:
        """
        批量导出多个任务
        
        Args:
            tasks: 任务列表，每个任务包含 task_id, description, result, steps, metadata
            formats: 导出格式列表，默认为 ['markdown']
            
        Returns:
            格式 -> 文件路径列表 的映射
        """
        if formats is None:
            formats = ['markdown']
        
        result = {fmt: [] for fmt in formats}
        
        for task in tasks:
            for fmt in formats:
                try:
                    _, file_path = self.export_task(
                        task_id=task.get('task_id', ''),
                        description=task.get('description', ''),
                        result=task.get('result', ''),
                        steps=task.get('steps', []),
                        format=fmt,
                        metadata=task.get('metadata'),
                        save_to_file=True
                    )
                    result[fmt].append(file_path)
                except Exception as e:
                    logger.error(f"Failed to export task {task.get('task_id')} to {fmt}: {e}")
        
        return result


# ============================================
# 通知系统
# ============================================

class NotificationManager:
    """通知管理器 - 支持多种通知方式"""
    
    def __init__(self):
        self.handlers = {}
    
    def register_handler(self, name: str, handler: callable):
        """注册通知处理器"""
        self.handlers[name] = handler
    
    async def notify(
        self,
        event: str,
        data: Dict,
        channels: List[str] = None
    ):
        """发送通知"""
        if channels is None:
            channels = list(self.handlers.keys())
        
        for channel in channels:
            if channel in self.handlers:
                try:
                    await self.handlers[channel](event, data)
                except Exception as e:
                    logger.error(f"Notification failed for {channel}: {e}")
    
    @staticmethod
    async def webhook_handler(url: str):
        """创建 Webhook 通知处理器"""
        async def handler(event: str, data: Dict):
            import aiohttp
            async with aiohttp.ClientSession() as session:
                payload = {
                    "event": event,
                    "data": data,
                    "timestamp": datetime.now().isoformat()
                }
                async with session.post(url, json=payload) as resp:
                    return resp.status == 200
        return handler
    
    @staticmethod
    async def email_handler(
        smtp_host: str,
        smtp_port: int,
        username: str,
        password: str,
        recipients: List[str]
    ):
        """创建邮件通知处理器"""
        async def handler(event: str, data: Dict):
            import smtplib
            from email.mime.text import MIMEText
            from email.mime.multipart import MIMEMultipart
            
            msg = MIMEMultipart()
            msg['Subject'] = f"JoinFlow 任务通知: {event}"
            msg['From'] = username
            msg['To'] = ', '.join(recipients)
            
            body = f"""
            事件: {event}
            时间: {datetime.now().isoformat()}
            
            详情:
            {json.dumps(data, indent=2, ensure_ascii=False)}
            """
            
            msg.attach(MIMEText(body, 'plain'))
            
            with smtplib.SMTP(smtp_host, smtp_port) as server:
                server.starttls()
                server.login(username, password)
                server.send_message(msg)
            
            return True
        return handler

